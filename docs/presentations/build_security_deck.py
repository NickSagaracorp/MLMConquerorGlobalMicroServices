"""
Genera la presentación de estado del trabajo de seguridad y 2FA de administración.

Hermano de build_deck.py: reutiliza su paleta y tipografía para que la serie se reconozca,
pero escribe en este repositorio (build_deck.py apunta su OUT_DIR a otro repo) y no arrastra
su estado a nivel de módulo.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------- Paleta: la misma de build_deck.py ----------------
NAVY      = RGBColor(0x14, 0x1B, 0x4D)
NAVY_DARK = RGBColor(0x0A, 0x10, 0x33)
NAVY_MID  = RGBColor(0x2A, 0x35, 0x6B)
PAPER     = RGBColor(0xF7, 0xF6, 0xF1)
GOLD      = RGBColor(0xC9, 0xA0, 0x3A)
GOLD_DEEP = RGBColor(0xA0, 0x7A, 0x1F)
CHARCOAL  = RGBColor(0x2B, 0x2D, 0x33)
MUTED     = RGBColor(0x6D, 0x72, 0x82)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_OK  = RGBColor(0x2E, 0x7D, 0x32)
RED_FLAG  = RGBColor(0xB3, 0x26, 0x1A)
AMBER     = RGBColor(0xB2, 0x7A, 0x0F)

HEADER_FONT = "Georgia"
BODY_FONT   = "Calibri"

OUT_PATH = Path(__file__).parent / "2026-08-28-seguridad-2fa-estado.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

FOOTER = "MLMConqueror Global Edition  ·  Seguridad y 2FA  ·  2026-08-28"


# ---------------- helpers ----------------
def add_rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    return sh


def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=14, bold=False,
             color=CHARCOAL, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def new_slide(page_num, *, dark=False):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY_DARK if dark else PAPER)
    if page_num:
        add_text(s, Inches(0.6), Inches(7.05), Inches(9), Inches(0.3), FOOTER,
                 size=9, color=MUTED if not dark else RGBColor(0x8A, 0x92, 0xB0))
        add_text(s, Inches(12.1), Inches(7.05), Inches(0.6), Inches(0.3), str(page_num),
                 size=9, color=MUTED if not dark else RGBColor(0x8A, 0x92, 0xB0),
                 align=PP_ALIGN.RIGHT)
    return s


def header(slide, kicker, title, dark=False):
    add_text(slide, Inches(0.6), Inches(0.45), Inches(9), Inches(0.3), kicker.upper(),
             size=10, bold=True, color=GOLD if dark else GOLD_DEEP)
    add_text(slide, Inches(0.6), Inches(0.75), Inches(12.1), Inches(0.8), title,
             font=HEADER_FONT, size=30, bold=True, color=WHITE if dark else NAVY)


def stat_card(slide, x, y, w, value, label, accent=NAVY):
    add_rect(slide, x, y, w, Inches(1.5), fill=WHITE, line=RGBColor(0xDD, 0xDC, 0xD5),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    add_text(slide, x, y + Inches(0.22), w, Inches(0.7), value,
             font=HEADER_FONT, size=34, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.15), y + Inches(0.98), w - Inches(0.3), Inches(0.4), label,
             size=11, color=MUTED, align=PP_ALIGN.CENTER)


def bullet_rows(slide, x, y, w, rows, *, gap=Inches(0.92), dot=NAVY):
    """rows = [(titulo, descripcion), ...] con un punto de color a la izquierda."""
    for i, (title, desc) in enumerate(rows):
        yy = y + i * gap
        add_rect(slide, x, yy + Inches(0.08), Inches(0.16), Inches(0.16),
                 fill=dot, shape=MSO_SHAPE.OVAL)
        add_text(slide, x + Inches(0.34), yy, w - Inches(0.34), Inches(0.3), title,
                 size=14, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.34), yy + Inches(0.30), w - Inches(0.34), Inches(0.6),
                 desc, size=11.5, color=CHARCOAL)


# ---------------- 1 · Portada ----------------
def slide_title():
    s = new_slide(None, dark=True)
    add_text(s, Inches(0.9), Inches(1.9), Inches(9), Inches(0.35),
             "ESTADO DE DESARROLLO  ·  SEGURIDAD", size=11, bold=True, color=GOLD)
    add_text(s, Inches(0.9), Inches(2.35), Inches(11), Inches(1.2),
             "Autenticación de dos factores", font=HEADER_FONT, size=46, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.5), Inches(10.5), Inches(0.9),
             "Prerrequisitos de seguridad, núcleo de verificación por tres canales,\n"
             "y login con dos factores en el portal de administración.",
             size=16, color=RGBColor(0xC8, 0xD0, 0xE8), spacing=6)
    add_rect(s, Inches(0.9), Inches(4.75), Inches(1.4), Inches(0.05), fill=GOLD)
    add_text(s, Inches(0.9), Inches(5.05), Inches(9), Inches(0.9),
             "2026-08-28  ·  20 commits en master  ·  1.565 pruebas automatizadas",
             size=13, color=RGBColor(0x9A, 0xA6, 0xC8))


# ---------------- 2 · Resumen ----------------
def slide_summary():
    s = new_slide(2)
    header(s, "resumen ejecutivo", "Qué se construyó")

    stat_card(s, Inches(0.6),  Inches(1.85), Inches(2.9), "3", "bugs de autenticación\nvivos, encontrados y cerrados", RED_FLAG)
    stat_card(s, Inches(3.75), Inches(1.85), Inches(2.9), "1.565", "pruebas automatizadas\nen verde", GREEN_OK)
    stat_card(s, Inches(6.90), Inches(1.85), Inches(2.9), "3", "canales de verificación\napp, correo y SMS", NAVY)
    stat_card(s, Inches(10.05), Inches(1.85), Inches(2.7), "0", "cambios visibles\npara los usuarios hoy", GOLD_DEEP)

    bullet_rows(s, Inches(0.6), Inches(3.75), Inches(12.1), [
        ("Prerrequisitos de seguridad",
         "La llave privada que firma todos los tokens estaba en el repositorio: cualquiera con acceso podía "
         "firmar un token de administrador sin contraseña. Rotada, y el sistema rechaza la anterior."),
        ("Núcleo de verificación",
         "Librería propia que emite y verifica códigos por aplicación de autenticación, correo y SMS, "
         "con límites de intentos y protección contra reutilización del mismo código."),
        ("Login con dos factores",
         "El administrador entra con código. Quien no lo tenga configurado queda en la pantalla de alta "
         "y no puede navegar hasta completarla."),
    ])


# ---------------- 3 · Los tres bugs ----------------
def slide_bugs():
    s = new_slide(3)
    header(s, "hallazgos", "Tres fallos que ya estaban ahí")
    add_text(s, Inches(0.6), Inches(1.62), Inches(12.1), Inches(0.4),
             "Ninguno se buscaba. Aparecieron al construir encima y ninguno era visible desde fuera.",
             size=12.5, color=MUTED)

    cards = [
        ("El código de verificación\nnunca pudo validarse",
         "El servicio rechazaba el mismo código que acababa de emitir, por un renombrado interno de datos. "
         "No se notó porque el correo tampoco se enviaba y la tabla de plantillas estaba vacía: "
         "tres fallos tapándose entre sí."),
        ("El bloqueo por intentos\nfallidos no existía",
         "Estaba configurado desde el principio —cinco intentos, quince minutos— pero el login nunca lo "
         "invocaba. En la práctica se podían probar contraseñas sin límite."),
        ("Dos servicios rechazaban\ntodo token legítimo",
         "Billing y CommissionEngine verificaban firmas con un tipo de llave que no corresponde. "
         "Sus endpoints devolvían 401 a cualquier petición autenticada; el negocio no se vio afectado "
         "porque nadie los consumía por HTTP."),
    ]
    x = Inches(0.6)
    for i, (title, body) in enumerate(cards):
        cx = x + i * Inches(4.15)
        add_rect(s, cx, Inches(2.2), Inches(3.85), Inches(3.15), fill=WHITE,
                 line=RGBColor(0xDD, 0xDC, 0xD5), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
        add_rect(s, cx + Inches(0.3), Inches(2.5), Inches(0.42), Inches(0.42),
                 fill=RED_FLAG, shape=MSO_SHAPE.OVAL)
        add_text(s, cx + Inches(0.3), Inches(2.56), Inches(0.42), Inches(0.3), str(i + 1),
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, cx + Inches(0.3), Inches(3.12), Inches(3.25), Inches(0.9), title,
                 font=HEADER_FONT, size=15, bold=True, color=NAVY, spacing=2)
        add_text(s, cx + Inches(0.3), Inches(4.05), Inches(3.25), Inches(1.5), body,
                 size=11, color=CHARCOAL)


# ---------------- 4 · Estado de activación ----------------
def slide_posture():
    s = new_slide(4)
    header(s, "estado", "Desplegado, pero apagado")
    add_text(s, Inches(0.6), Inches(1.62), Inches(12.1), Inches(0.5),
             "Todo el código está en master y no cambia nada para nadie todavía. "
             "Encenderlo tiene un orden, y saltárselo deja gente fuera del portal.",
             size=13, color=CHARCOAL)

    rows = [
        ("Hoy", "Enrolamiento obligatorio desactivado. Los códigos se escriben en el registro, no se envían.", GREEN_OK),
        ("Paso 1", "Configurar las credenciales de correo y SMS en el entorno de producción.", NAVY),
        ("Paso 2", "Verificar que un correo y un mensaje llegan de verdad a una cuenta real.", NAVY),
        ("Paso 3", "Rotar la llave de firma en producción, en ventana de mantenimiento.", NAVY),
        ("Paso 4", "Activar el enrolamiento obligatorio y avisar al equipo de administración.", GOLD_DEEP),
    ]
    for i, (tag, desc, color) in enumerate(rows):
        yy = Inches(2.5) + i * Inches(0.86)
        add_rect(s, Inches(0.6), yy, Inches(1.25), Inches(0.55), fill=color,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
        add_text(s, Inches(0.6), yy + Inches(0.13), Inches(1.25), Inches(0.3), tag,
                 size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.1), yy + Inches(0.12), Inches(10.6), Inches(0.5), desc,
                 size=13, color=CHARCOAL)



# ---------------- 5 · Riesgos ----------------
def slide_risks():
    s = new_slide(5)
    header(s, "riesgos abiertos", "Lo que todavía no está resuelto")

    rows = [
        ("La llave anterior sigue en el historial del repositorio",
         "Se rotó y el sistema la rechaza, pero cualquiera con una copia antigua la tiene. "
         "Borrarla del historial exige reescribirlo y coordinar con todas las copias existentes.", RED_FLAG),
        ("La entrega real de correos y mensajes no se ha probado",
         "Las pruebas automatizadas aíslan la llamada al proveedor. Que un mensaje llegue de verdad "
         "solo se comprueba con credenciales reales.", AMBER),
        ("Las traducciones no están revisadas",
         "Los textos de las pantallas nuevas se tradujeron a los nueve idiomas siguiendo la terminología "
         "existente, pero ningún hablante nativo los ha revisado.", AMBER),
        ("El cambio de horario de verano",
         "El sistema pasó a usar hora de servidor. En el cambio de otoño la misma hora ocurre dos veces, "
         "lo que afecta a lo que se ordena por fecha. No se ha abordado.", AMBER),
    ]
    for i, (title, desc, color) in enumerate(rows):
        yy = Inches(1.95) + i * Inches(1.24)
        add_rect(s, Inches(0.6), yy, Inches(0.08), Inches(1.0), fill=color)
        add_text(s, Inches(0.95), yy, Inches(11.8), Inches(0.35), title,
                 size=14.5, bold=True, color=NAVY)
        add_text(s, Inches(0.95), yy + Inches(0.36), Inches(11.8), Inches(0.65), desc,
                 size=11.5, color=CHARCOAL)


# ---------------- 6 · Lo que falta ----------------
def slide_next():
    s = new_slide(6)
    header(s, "siguiente", "Lo que queda por construir")

    # Altura ajustada al contenido: con 4.3" quedaba medio panel vacío al pie.
    CARD_H = Inches(3.85)

    add_rect(s, Inches(0.6), Inches(1.9), Inches(6.0), CARD_H, fill=WHITE,
             line=RGBColor(0xDD, 0xDC, 0xD5), shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    add_text(s, Inches(1.0), Inches(2.25), Inches(5.2), Inches(0.85),
             "Superficie de cuenta", font=HEADER_FONT, size=19, bold=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(2.95), Inches(5.2), Inches(2.6),
             "Recuperación de contraseña.\n"
             "Confirmación de correo.\n"
             "Alta y verificación de teléfono.\n"
             "Gestión del segundo factor desde el perfil.\n"
             "Datos personales y descarga.\n\n"
             "Hoy ninguna de estas pantallas existe en el portal de administración.",
             size=13, color=CHARCOAL, spacing=7)

    add_rect(s, Inches(7.0), Inches(1.9), Inches(5.7), CARD_H, fill=NAVY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    # 0.85" de alto: este título envuelve a dos líneas y con 0.4" pisaba el cuerpo.
    add_text(s, Inches(7.4), Inches(2.25), Inches(4.9), Inches(0.85),
             "Confirmación en operaciones críticas", font=HEADER_FONT, size=19, bold=True, color=WHITE)
    add_text(s, Inches(7.4), Inches(3.25), Inches(4.9), Inches(2.3),
             "Pedir un código antes de liberar un pago, cambiar un rol, tocar credenciales "
             "de cobro o ajustar datos de negocio de alto impacto.\n\n"
             "Qué operación lo exige y por qué canal se configura sin desplegar.\n\n"
             "Era la segunda mitad del encargo original y todavía no está construida.",
             size=13, color=RGBColor(0xC8, 0xD0, 0xE8), spacing=7)


# ---------------- 7 · Cierre ----------------
def slide_closing():
    s = new_slide(None, dark=True)
    add_text(s, Inches(0.9), Inches(2.6), Inches(11), Inches(1.0),
             "Preguntas", font=HEADER_FONT, size=44, bold=True, color=WHITE)
    add_rect(s, Inches(0.9), Inches(3.75), Inches(1.4), Inches(0.05), fill=GOLD)
    add_text(s, Inches(0.9), Inches(4.1), Inches(10.5), Inches(1.2),
             "Documentación de referencia:\n"
             "docs/deployment/verificacion-2fa.pdf  ·  qué comprobar y en qué orden\n"
             "docs/deployment/rotacion-llaves-jwt.pdf  ·  cómo rotar la llave en producción",
             size=14, color=RGBColor(0xC8, 0xD0, 0xE8), spacing=6)


# ---------------- Build ----------------
if __name__ == "__main__":
    slide_title()
    slide_summary()
    slide_bugs()
    slide_posture()
    slide_risks()
    slide_next()
    slide_closing()

    prs.save(OUT_PATH)
    print(f"Guardado: {OUT_PATH}")
    print(f"Diapositivas: {len(prs.slides)}")
    print(f"Tamaño: {OUT_PATH.stat().st_size / 1024:.1f} KB")
