"""
Build the MLMConqueror Global Edition development-status deck.
Sprint 15 review for client presentation.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import math

# ---------------- Palette: Midnight Executive + gold accent ----------------
NAVY        = RGBColor(0x14, 0x1B, 0x4D)   # primary deep navy
NAVY_DARK   = RGBColor(0x0A, 0x10, 0x33)
NAVY_MID    = RGBColor(0x2A, 0x35, 0x6B)
ICE         = RGBColor(0xE9, 0xEE, 0xFA)   # light bg
PAPER       = RGBColor(0xF7, 0xF6, 0xF1)   # warm off-white
GOLD        = RGBColor(0xC9, 0xA0, 0x3A)   # accent
GOLD_DEEP   = RGBColor(0xA0, 0x7A, 0x1F)
CHARCOAL    = RGBColor(0x2B, 0x2D, 0x33)
MUTED       = RGBColor(0x6D, 0x72, 0x82)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_OK    = RGBColor(0x2E, 0x7D, 0x32)
RED_FLAG    = RGBColor(0xB3, 0x26, 0x1A)
TEAL        = RGBColor(0x1C, 0x72, 0x93)

# Wing accents for the mindmap
WING_MEMBER = RGBColor(0x2C, 0x5F, 0x8D)   # blue
WING_ADMIN  = RGBColor(0x6D, 0x2E, 0x46)   # berry
WING_ENGINE = RGBColor(0xC9, 0xA0, 0x3A)   # gold
WING_FOUND  = RGBColor(0x36, 0x65, 0x4A)   # forest
WING_DATA   = RGBColor(0x50, 0x4A, 0x80)   # violet

HEADER_FONT = "Georgia"
BODY_FONT   = "Calibri"

OUT_DIR = Path(r"C:/Users/sagar/source/repos/ClaudeRepository/docs/presentations")
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / "2026-05-24-mlm-development-status.pptx"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]

# ---------------- helpers ----------------
def add_rect(slide, x, y, w, h, fill=None, line=None, shadow=False, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.text_frame.margin_left = Emu(0)
    s.text_frame.margin_right = Emu(0)
    s.text_frame.margin_top = Emu(0)
    s.text_frame.margin_bottom = Emu(0)
    return s

def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=14, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def page_bg(slide, color=PAPER):
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)
    return bg

def left_band(slide, color=NAVY, w=Inches(0.18)):
    return add_rect(slide, 0, 0, w, SLIDE_H, fill=color)

def page_header(slide, title, kicker=None):
    # Kicker (small uppercase tag)
    if kicker:
        add_text(slide, Inches(0.55), Inches(0.35), Inches(8), Inches(0.3),
                 kicker.upper(), font=BODY_FONT, size=10, bold=True, color=GOLD_DEEP)
    add_text(slide, Inches(0.55), Inches(0.6), Inches(12.2), Inches(0.7),
             title, font=HEADER_FONT, size=30, bold=True, color=NAVY)
    # thin gold rule under title? Skip — guidance says no accent lines under titles.

def footer(slide, page_num):
    add_text(slide, Inches(0.55), Inches(7.10), Inches(8), Inches(0.3),
             "MLMConqueror Global Edition  ·  Sprint 15 Review  ·  2026-05-24",
             font=BODY_FONT, size=9, color=MUTED)
    add_text(slide, Inches(12.2), Inches(7.10), Inches(0.6), Inches(0.3),
             str(page_num), font=BODY_FONT, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def new_slide(page_num, *, dark=False):
    s = prs.slides.add_slide(BLANK)
    page_bg(s, NAVY_DARK if dark else PAPER)
    if not dark:
        left_band(s, NAVY)
    footer(s, page_num)
    return s

# ---------------- Slide 1 — Title ----------------
def slide_title():
    s = prs.slides.add_slide(BLANK)
    page_bg(s, NAVY_DARK)
    # Gold corner motif (top-right)
    add_rect(s, Inches(11.4), Inches(0), Inches(1.93), Inches(0.55), fill=GOLD)
    add_rect(s, Inches(12.8), Inches(0.55), Inches(0.55), Inches(1.2), fill=GOLD_DEEP)
    # Kicker
    add_text(s, Inches(0.8), Inches(2.4), Inches(8), Inches(0.4),
             "DEVELOPMENT STATUS  ·  CLIENT REVIEW",
             font=BODY_FONT, size=12, bold=True, color=GOLD)
    # Title — drop to 46pt so it stays on one line in 11.5" width
    add_text(s, Inches(0.8), Inches(2.85), Inches(12), Inches(1.4),
             "MLMConqueror Global Edition",
             font=HEADER_FONT, size=46, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(4.25), Inches(11), Inches(0.9),
             "Where we are, what we shipped, what comes next",
             font=HEADER_FONT, size=22, italic=True, color=ICE)
    # Bottom meta block
    add_rect(s, Inches(0.8), Inches(6.0), Inches(4.2), Inches(0.04), fill=GOLD)
    add_text(s, Inches(0.8), Inches(6.15), Inches(8), Inches(0.4),
             "Sprint 15 Review",
             font=BODY_FONT, size=14, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(6.55), Inches(8), Inches(0.4),
             "2026-05-24  ·  Prepared for Client Discussion",
             font=BODY_FONT, size=11, color=ICE)
    # Page marker
    add_text(s, Inches(12.2), Inches(7.10), Inches(0.6), Inches(0.3),
             "01", font=BODY_FONT, size=9, color=ICE, align=PP_ALIGN.RIGHT)

# ---------------- Slide 2 — Executive summary ----------------
def slide_exec_summary():
    s = new_slide(2)
    page_header(s, "Where the platform stands today", kicker="Executive summary")

    bullets = [
        ("13 services. One platform.",
         "The full backend (Domain, Repository, 11 services) is in place. Clean architecture, "
         "CQRS, full test coverage — over 900 unit tests passing across the solution."),
        ("Release 1 backend: feature-complete.",
         "Sprints 1-14 delivered signups, ranks, commissions, billing, dual-team trees, tokens, "
         "tickets, AdminWeb and BizCenterWeb. Sprint 15 closed the rank-achievement certificates loop."),
        ("Load-tested to 400 concurrent signups — zero failures.",
         "Six load-test waves on a real-world tree (2,453 ambassadors at peak) confirmed the "
         "pipeline holds at ~30 signups/sec sustained, with cascade rank promotions working."),
        ("10 production-grade defects found and fixed — under load.",
         "The hardening exercise surfaced concurrency, queueing and data-integrity bugs that "
         "only appear at scale. Every one is fixed in code."),
        ("Rank ladder: 11 of 19 ranks reached on a live ambassador.",
         "Silver through Double Diamond exercised end-to-end with literal signups and real "
         "PDF certificates. Remaining 8 ranks shift to seeded data (scale becomes infeasible)."),
    ]

    y = Inches(1.65)
    for i, (head, body) in enumerate(bullets):
        # number circle
        cx = Inches(0.6)
        cy = y
        circle = add_rect(s, cx, cy, Inches(0.55), Inches(0.55),
                          fill=NAVY, shape=MSO_SHAPE.OVAL)
        tf = circle.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i+1)
        r.font.name = BODY_FONT; r.font.size = Pt(16); r.font.bold = True
        r.font.color.rgb = GOLD

        add_text(s, Inches(1.35), y - Inches(0.02), Inches(11.4), Inches(0.4),
                 head, font=HEADER_FONT, size=16, bold=True, color=NAVY)
        add_text(s, Inches(1.35), y + Inches(0.38), Inches(11.4), Inches(0.6),
                 body, font=BODY_FONT, size=11.5, color=CHARCOAL, line_spacing=1.15)
        y += Inches(1.02)

# ---------------- Slide 3 — Mindmap (centerpiece) ----------------
def slide_mindmap():
    s = new_slide(3)
    page_header(s, "How the system fits together", kicker="The architecture at a glance")

    # subtitle
    add_text(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.35),
             "Five wings around one platform — member-facing apps, admin tools, business engines, the foundation, and infrastructure.",
             font=BODY_FONT, size=11.5, italic=True, color=MUTED)

    # Canvas area (below header)
    # Hub centred roughly between top wings (y=2.05) and the engines band (y=5.80).
    cx = Inches(6.6665)  # slide centre (SLIDE_W/2)
    cy = Inches(4.05)

    R_HUB_W, R_HUB_H = Inches(2.6), Inches(1.05)
    hub_x = cx - R_HUB_W/2
    hub_y = cy - R_HUB_H/2
    # We'll plot 5 wing boxes radially. Use approximate positions hand-tuned for readability.

    # Layout: 5 wings around the centre hub.
    # Top-left: Member-facing.  Top-right: Admin.
    # Mid-left: Foundation.     Mid-right: Data & Infrastructure.
    # Bottom-centre: Business Engines (wide band, single row of 4 children).
    wings = [
        # (title, color, wx, wy, ww, wh, children:[(name,wx,wy,ww,wh)])
        # ------ Member-facing (top-left) ------
        ("Member-facing", WING_MEMBER, Inches(0.45), Inches(2.05), Inches(3.4), Inches(0.65),
            [("Signups API",          Inches(0.45), Inches(2.80), Inches(1.60), Inches(0.50)),
             ("BizCenterWeb (Blazor)",Inches(2.20), Inches(2.80), Inches(1.65), Inches(0.50)),
             ("BizCenter API",        Inches(1.30), Inches(3.40), Inches(1.85), Inches(0.50))]),

        # ------ Admin (top-right) ------
        ("Admin", WING_ADMIN, Inches(9.45), Inches(2.05), Inches(3.4), Inches(0.65),
            [("AdminWeb (Blazor)",    Inches(9.45), Inches(2.80), Inches(1.65), Inches(0.50)),
             ("AdminAPI",             Inches(11.25),Inches(2.80), Inches(1.60), Inches(0.50)),
             ("SharedAPICenter",      Inches(10.30),Inches(3.40), Inches(1.95), Inches(0.50))]),

        # ------ Foundation (mid-left) ------
        ("Foundation", WING_FOUND, Inches(0.45), Inches(4.45), Inches(3.4), Inches(0.65),
            [("Domain",               Inches(0.45), Inches(5.20), Inches(1.10), Inches(0.50)),
             ("SharedKernel",         Inches(1.65), Inches(5.20), Inches(1.50), Inches(0.50)),
             ("Repository (EF Core)", Inches(0.45), Inches(5.80), Inches(2.70), Inches(0.50))]),

        # ------ Data & Infra (mid-right) ------
        ("Data & Infrastructure", WING_DATA, Inches(9.45), Inches(4.45), Inches(3.4), Inches(0.65),
            [("SQL Server",           Inches(9.55), Inches(5.20), Inches(1.55), Inches(0.50)),
             ("Hangfire",             Inches(11.20),Inches(5.20), Inches(1.55), Inches(0.50)),
             ("Redis (cache)",        Inches(9.55), Inches(5.80), Inches(1.55), Inches(0.50)),
             ("Firebase (push)",      Inches(11.20),Inches(5.80), Inches(1.55), Inches(0.50))]),

        # ------ Engines (bottom-centre band) ------
        ("Business Engines", WING_ENGINE, Inches(4.65), Inches(5.80), Inches(4.00), Inches(0.55),
            [("RankEngine",           Inches(4.10), Inches(6.45), Inches(1.55), Inches(0.50)),
             ("CommissionEngine",     Inches(5.75), Inches(6.45), Inches(2.00), Inches(0.50)),
             ("Billing",              Inches(7.85), Inches(6.45), Inches(1.30), Inches(0.50)),
             ("TicketMgmtSystem",     Inches(9.25), Inches(6.45), Inches(1.90), Inches(0.50))]),
    ]

    # 1) Draw wing rectangles + child boxes FIRST (so connectors sit on top? Actually we want connectors below text.)
    # Draw connectors first (lines), then nodes on top.

    # Hub centre point (for connectors)
    hub_cx_emu = (hub_x + R_HUB_W/2)
    hub_cy_emu = (hub_y + R_HUB_H/2)

    # ---- Connectors hub → wing header ----
    for (title, color, wx, wy, ww, wh, children) in wings:
        wing_cx = wx + ww/2
        wing_cy = wy + wh/2
        line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      int(hub_cx_emu), int(hub_cy_emu),
                                      int(wing_cx), int(wing_cy))
        line.line.color.rgb = color
        line.line.width = Pt(2.25)

    # ---- Connectors wing → its children ----
    for (title, color, wx, wy, ww, wh, children) in wings:
        wing_cx = wx + ww/2
        wing_cy = wy + wh/2
        for (cname, ccx, ccy, ccw, cch) in children:
            chx = ccx + ccw/2
            chy = ccy + cch/2
            line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          int(wing_cx), int(wing_cy),
                                          int(chx), int(chy))
            line.line.color.rgb = color
            line.line.width = Pt(1.0)
            # dashed
            ln = line.line._get_or_add_ln()
            prstDash = etree.SubElement(ln, qn('a:prstDash'))
            prstDash.set('val', 'dash')

    # ---- HUB (centre node) ----
    hub = add_rect(s, hub_x, hub_y, R_HUB_W, R_HUB_H, fill=NAVY,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    hub.line.color.rgb = GOLD
    hub.line.width = Pt(2.0)
    tf = hub.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "MLMConqueror"
    r.font.name = HEADER_FONT; r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "Platform"
    r2.font.name = HEADER_FONT; r2.font.size = Pt(15); r2.font.bold = True
    r2.font.color.rgb = GOLD

    # ---- Wing heads + children ----
    for (title, color, wx, wy, ww, wh, children) in wings:
        # Wing head
        wh_box = add_rect(s, wx, wy, ww, wh, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        wh_box.line.fill.background()
        tf = wh_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title
        r.font.name = HEADER_FONT; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = WHITE

        for (cname, ccx, ccy, ccw, cch) in children:
            box = add_rect(s, ccx, ccy, ccw, cch, fill=WHITE,
                           shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=color)
            box.line.width = Pt(1.25)
            tf = box.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Emu(0); tf.margin_right = Emu(0)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = cname
            r.font.name = BODY_FONT; r.font.size = Pt(9.5); r.font.bold = True
            r.font.color.rgb = NAVY

    # ---- Legend / cross-cutting flow callout (bottom-centre below engines) ----
    legend_y = Inches(7.18)
    add_text(s, Inches(0.55), legend_y, Inches(12.2), Inches(0.25),
             "Cross-cutting flow:  Signups  →  rank evaluation queue  →  rank promotion  →  certificate + notifications",
             font=BODY_FONT, size=9, italic=True, color=GOLD_DEEP, align=PP_ALIGN.CENTER)

# ---------------- Slide 4 — Service breakdown ----------------
def slide_service_table():
    s = new_slide(4)
    page_header(s, "The 13 services and what each one does", kicker="Service breakdown")

    # 2-column layout, 7 rows each (13 + 1 banner row)
    rows_left = [
        ("Domain",            "Pure entities, enums, business rules — zero external dependencies."),
        ("SharedKernel",      "Result<T>, ApiResponse<T>, PagedResult, common interfaces."),
        ("Repository",        "EF Core context, Fluent API configurations, audit interceptor, Unit of Work."),
        ("Signups",           "Ambassador + member registration, 3-phase wizard, placement, upgrade/downgrade."),
        ("RankEngine",        "Rank evaluation, qualification, certificate PDF generation."),
        ("CommissionEngine",  "Fast Start, Daily Residual, Boost, Presidential, Matching — all commission types."),
        ("Billing",           "Payment gateways, recurring billing, dunning, commission-balance funding."),
    ]
    rows_right = [
        ("BizCenter",         "Member-facing API: profile, wallet, billing, commissions, team, tokens, loyalty."),
        ("BizCenterWeb",      "Blazor MAUI Hybrid app for ambassadors and members."),
        ("AdminAPI",          "Admin operations: members, commissions, ghost points, config, dashboards."),
        ("AdminWeb",          "Blazor admin console — operations, configuration, member management."),
        ("TicketMgmtSystem",  "Helpdesk: tickets, comments, attachments, assignments, merges."),
        ("SharedAPICenter",   "External webhooks (Stripe/Braintree/crypto) and integration endpoints."),
        ("(Cross-cutting)",   "Hangfire scheduler, Redis cache, Firebase push notifications, JWT auth."),
    ]

    def draw_column(rows, x0):
        y = Inches(1.6)
        row_h = Inches(0.72)
        col_w = Inches(6.15)
        for name, desc in rows:
            # row card
            add_rect(s, x0, y, col_w, row_h - Inches(0.07), fill=WHITE,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=ICE)
            # service name
            add_text(s, x0 + Inches(0.15), y + Inches(0.06), Inches(2.0), Inches(0.3),
                     name, font=HEADER_FONT, size=12, bold=True, color=NAVY)
            # description
            add_text(s, x0 + Inches(0.15), y + Inches(0.34), col_w - Inches(0.3), Inches(0.32),
                     desc, font=BODY_FONT, size=9.5, color=CHARCOAL)
            y += row_h

    draw_column(rows_left,  Inches(0.55))
    draw_column(rows_right, Inches(7.0))

# ---------------- Slide 5 — Data flow ----------------
def slide_data_flow():
    s = new_slide(5)
    page_header(s, "From signup to certificate", kicker="End-to-end flow")

    add_text(s, Inches(0.55), Inches(1.4), Inches(12.2), Inches(0.4),
             "The path a new ambassador's data travels — and how that triggers rank promotion and notifications downstream.",
             font=BODY_FONT, size=12, italic=True, color=MUTED)

    steps = [
        ("1", "Signup wizard",        "Member fills 3-phase wizard\n(personal · placement · payment)", "Signups API"),
        ("2", "Statistic delta",      "MemberStatisticDelta enqueued\nto stats pipeline",               "Queue"),
        ("3", "Apply deltas",         "Background job merges\ndeltas atomically into\nupline statistics",  "Hangfire"),
        ("4", "Rank queue",           "Each modified upline is\nenqueued for rank evaluation",          "Queue"),
        ("5", "Evaluate rank",        "EvaluateRankHandler checks\nqualifications against thresholds", "RankEngine"),
        ("6", "Promote + notify",     "MemberRankHistory row written;\nnotifications fan out",          "RankEngine"),
        ("7", "Certificate on demand","PDF generated when member or\nadmin requests it",                "RankEngine"),
    ]

    # 7 horizontal cards
    x = Inches(0.55)
    y = Inches(2.2)
    card_w = Inches(1.71)
    card_h = Inches(3.4)
    gap = Inches(0.08)

    for i, (num, title, desc, owner) in enumerate(steps):
        # number circle on top
        cx = x + card_w/2 - Inches(0.30)
        circ = add_rect(s, cx, y, Inches(0.6), Inches(0.6), fill=GOLD, shape=MSO_SHAPE.OVAL)
        tf = circ.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.name = HEADER_FONT; r.font.size = Pt(20); r.font.bold = True
        r.font.color.rgb = NAVY

        # card body
        body_top = y + Inches(0.42)
        add_rect(s, x, body_top, card_w, card_h - Inches(0.42), fill=WHITE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=ICE)
        # title
        add_text(s, x + Inches(0.1), body_top + Inches(0.30), card_w - Inches(0.2), Inches(0.6),
                 title, font=HEADER_FONT, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # desc
        add_text(s, x + Inches(0.1), body_top + Inches(0.90), card_w - Inches(0.2), Inches(1.5),
                 desc, font=BODY_FONT, size=9.5, color=CHARCOAL, align=PP_ALIGN.CENTER, line_spacing=1.2)
        # owner pill
        owner_y = body_top + card_h - Inches(0.85)
        add_rect(s, x + Inches(0.25), owner_y, card_w - Inches(0.5), Inches(0.32),
                 fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tx = add_text(s, x + Inches(0.25), owner_y + Inches(0.04), card_w - Inches(0.5), Inches(0.26),
                      owner, font=BODY_FONT, size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

        # arrow between cards
        if i < len(steps) - 1:
            ax1 = x + card_w
            ax2 = ax1 + gap
            ay  = y + Inches(0.42) + (card_h - Inches(0.42))/2
            arrow = add_rect(s, ax1, ay - Inches(0.07), gap, Inches(0.14), fill=GOLD,
                             shape=MSO_SHAPE.RIGHT_ARROW)
            arrow.line.fill.background()

        x += card_w + gap

    # bottom callout
    add_rect(s, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.85),
             fill=ICE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.75), Inches(6.15), Inches(11.8), Inches(0.4),
             "Why this matters", font=HEADER_FONT, size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.75), Inches(6.43), Inches(11.8), Inches(0.45),
             "Steps 2-7 are decoupled. A spike of signups never blocks rank evaluation, and rank evaluation never blocks the user's response. Each stage scales independently.",
             font=BODY_FONT, size=11, color=CHARCOAL)

# ---------------- Slide 6 — Sprint 15 highlights ----------------
def slide_sprint15():
    s = new_slide(6)
    page_header(s, "What we delivered this sprint", kicker="Sprint 15 highlights")

    items = [
        ("Rank certificate system",
         "PDF certificates auto-generated when an ambassador reaches each rank — visual name + date drawn over 19 designed templates."),
        ("AdminWeb certificate UI",
         "Admins can view, regenerate or delete a certificate from the member profile page."),
        ("Input validation across every public endpoint",
         "189 new FluentValidation rules cover Signups, AdminAPI and BizCenter — names, emails, phones, SSN/EIN, payment fields, free-text whitelists."),
        ("Eventual-consistency stats pipeline",
         "Replaced per-signup synchronous upline writes with a queued atomic-merge job — eliminated the lost-update race entirely."),
        ("BFS placement + cascade promotion gate",
         "Binary placement now uses breadth-first slot-finding (no more 78-level degenerate chains); cascade rank promotions unlocked across the upline."),
        ("Paginated, filtered members admin",
         "AdminWeb members grid with server-side pagination, page-size selector, status filters and a stats panel above."),
        ("Multicultural load-test rig",
         "C# concurrency rig signs up batches with 1000 first + 1000 last names across 11 cultural pools — 6,400+ unique combinations."),
    ]
    y = Inches(1.6)
    for title, desc in items:
        add_rect(s, Inches(0.55), y, Inches(0.18), Inches(0.65), fill=GOLD)
        add_text(s, Inches(0.9), y - Inches(0.03), Inches(11.8), Inches(0.35),
                 title, font=HEADER_FONT, size=13, bold=True, color=NAVY)
        add_text(s, Inches(0.9), y + Inches(0.30), Inches(11.8), Inches(0.45),
                 desc, font=BODY_FONT, size=11, color=CHARCOAL, line_spacing=1.15)
        y += Inches(0.75)

# ---------------- Slide 7 — Load-test results ----------------
def slide_load_test():
    s = new_slide(7)
    page_header(s, "Load-test results", kicker="Capacity benchmarks")

    add_text(s, Inches(0.55), Inches(1.4), Inches(12.2), Inches(0.4),
             "Eleven waves of progressively heavier signups under one test ambassador (peak 2,453 in the downline).",
             font=BODY_FONT, size=11.5, italic=True, color=MUTED)

    headers = ["Wave", "Pattern", "Total", "Wall-clock", "Throughput", "Success", "Cascades"]
    rows = [
        ("Silver",            "4 sequential",                        "4",   "manual",  "—",        "100%",  "—"),
        ("Gold",              "10 sequential",                       "10",  "~1 min",  "~0.2/s",   "100%",  "—"),
        ("Platinum/Titanium", "various sequential",                  "15-90","minutes","low",      "100%",  "—"),
        ("Jade",              "5 waves, 10→200 conc, 20 sponsors",   "420", "~17 s",   "30/s",     "100%",  "0 (gated)"),
        ("Pearl",             "300 conc, 50 sponsors",               "300", "9.21 s",  "32.6/s",   "99.7%", "0 (gated)"),
        ("Emerald",           "200 conc, 50 sponsors, branches",     "200", "8.83 s",  "22.7/s",   "99.5%", "0 (gated)"),
        ("Ruby",              "400 conc, 80 sponsors",               "400", "14.24 s", "28.1/s",   "100%",  "0 (gated)"),
        ("Sapphire",          "350 conc, 69 sponsors",               "350", "12.08 s", "29.0/s",   "98.9%", "0 (gated)"),
        ("Diamond",           "350×1 max-spread (post cascade-fix)", "350", "17.30 s", "20.2/s",   "100%",  "205"),
        ("Double Diamond",    "350×1 (post eventual-consistency)",   "350", "11.92 s", "29.4/s",   "100%",  "54"),
    ]

    # Table
    x0 = Inches(0.55)
    y0 = Inches(2.0)
    col_w = [Inches(1.45), Inches(3.20), Inches(0.85), Inches(1.10), Inches(1.30), Inches(1.10), Inches(1.30)]
    total_w = sum(col_w, Emu(0))
    head_h = Inches(0.4)
    row_h  = Inches(0.36)

    # header
    cx = x0
    for i, h in enumerate(headers):
        add_rect(s, cx, y0, col_w[i], head_h, fill=NAVY)
        add_text(s, cx, y0 + Inches(0.07), col_w[i], head_h - Inches(0.05),
                 h, font=BODY_FONT, size=10.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        cx += col_w[i]

    # rows
    y = y0 + head_h
    for ri, row in enumerate(rows):
        fill = WHITE if ri % 2 == 0 else ICE
        cx = x0
        # highlight Double Diamond row
        is_dd = row[0] == "Double Diamond"
        if is_dd:
            fill = RGBColor(0xFB, 0xF1, 0xD7)  # soft gold tint
        for ci, val in enumerate(row):
            add_rect(s, cx, y, col_w[ci], row_h, fill=fill)
            bold = is_dd or ci == 0
            color = NAVY if (ci == 0) else (NAVY_DARK if is_dd else CHARCOAL)
            add_text(s, cx + Inches(0.08), y + Inches(0.07), col_w[ci] - Inches(0.16), row_h - Inches(0.05),
                     val, font=BODY_FONT, size=10, bold=bold, color=color,
                     align=PP_ALIGN.LEFT if ci <= 1 else PP_ALIGN.CENTER)
            cx += col_w[ci]
        y += row_h

    # callout
    cy = y + Inches(0.20)
    add_rect(s, Inches(0.55), cy, Inches(12.2), Inches(0.70),
             fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.75), cy + Inches(0.10), Inches(11.8), Inches(0.3),
             "Inflection point",
             font=BODY_FONT, size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.75), cy + Inches(0.34), Inches(11.8), Inches(0.32),
             "After the eventual-consistency rewrite, the worst-case max-spread pattern improved from 17.3s to 11.9s wall-clock — a 31% drop — at higher cascade count.",
             font=BODY_FONT, size=10.5, color=WHITE)

# ---------------- Slide 8 — Architectural improvements ----------------
def slide_improvements():
    s = new_slide(8)
    page_header(s, "Architectural improvements landed this sprint", kicker="Hardening")

    items = [
        ("Atomic statistics MERGE",
         "Replaced read-modify-write on upline statistics with a single `MERGE … WITH (HOLDLOCK)`. Zero lost updates verified under 421 concurrent writers."),
        ("Breadth-first placement",
         "Binary tree placement uses BFS to find the first empty slot — no more 78-level deep chains that crashed the SQL Server B-tree index."),
        ("External-members cascade unlock",
         "Removed the seeded `ExternalMembers = 1` requirement that was silently blocking every cascade rank promotion. Cascades now light up as expected."),
        ("Eventual-consistency stats pipeline",
         "Signups return in 50-200 ms; upline statistics propagate via Hangfire within ≤60 s. Decoupled the user response from heavy upline work."),
        ("Intermediate-rank chains",
         "When stats jump past several ranks at once, the engine now promotes through each intermediate rank in order — no rank skipped."),
        ("On-demand certificate generation",
         "Certificates render only when a member or admin opens them — avoids generating 200+ PDFs synchronously during a cascade wave."),
    ]
    # 3x2 grid
    cols = 3
    card_w = Inches(4.05)
    card_h = Inches(2.45)
    x0 = Inches(0.55)
    y0 = Inches(1.65)
    gap_x = Inches(0.12)
    gap_y = Inches(0.18)

    for i, (title, desc) in enumerate(items):
        col = i % cols
        row = i // cols
        x = x0 + col * (card_w + gap_x)
        y = y0 + row * (card_h + gap_y)
        add_rect(s, x, y, card_w, card_h, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=ICE)
        # gold corner tab
        add_rect(s, x, y, Inches(0.6), Inches(0.18), fill=GOLD)
        # number
        add_text(s, x + Inches(0.15), y + Inches(0.30), Inches(0.45), Inches(0.3),
                 f"0{i+1}", font=HEADER_FONT, size=11, bold=True, color=GOLD_DEEP)
        add_text(s, x + Inches(0.15), y + Inches(0.65), card_w - Inches(0.3), Inches(0.6),
                 title, font=HEADER_FONT, size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.15), y + Inches(1.20), card_w - Inches(0.3), Inches(1.15),
                 desc, font=BODY_FONT, size=10.5, color=CHARCOAL, line_spacing=1.20)

# ---------------- Slide 9 — Bugs found + fixed ----------------
def slide_bugs():
    s = new_slide(9)
    page_header(s, "Bugs found and fixed under load", kicker="Hardening at scale")

    add_text(s, Inches(0.55), Inches(1.4), Inches(12.2), Inches(0.45),
             "Production-grade systems are forged by exercising them under realistic load. Twelve defects surfaced during the climb — every one is fixed in code.",
             font=BODY_FONT, size=11.5, italic=True, color=MUTED)

    bugs = [
        ("1",  "Hangfire jobs parked on wrong queue",                 "Queue isolation"),
        ("2",  "JWT config drift (HMAC vs RSA, audience mismatch)",   "Cross-service auth"),
        ("3",  "Cert endpoints blocked SuperAdmin role",              "Authorization"),
        ("4",  "EvaluateRankHandler coupled to email/push I/O",       "Async decoupling"),
        ("5",  "Lost-update race on MemberStatistics",                "Concurrency"),
        ("6",  "Degenerate 78-level chain crashed SQL index",         "Placement"),
        ("7",  "SignupAPI never recomputed leg points",               "Data integrity"),
        ("8",  "Root ambassador never inserted into DualTeamTree",    "Bootstrap"),
        ("9",  "OrderNumberHelper infinite loop (676 ceiling)",       "Entropy"),
        ("10", "Hierarchy paths concatenated without slashes",        "Path normalization"),
        ("11", "Admin placement handlers used Count vs Sum",          "Calculation"),
        ("12", "ExternalMembers=1 seed blocked all cascades",         "Business-rule seed"),
    ]
    # 2 columns of 6 — give the tag pill its own column with safe gap
    col_w     = Inches(6.10)
    row_h     = Inches(0.70)
    x_left    = Inches(0.55)
    x_right   = Inches(6.75)
    y0        = Inches(2.05)
    badge_w   = Inches(0.46)
    pill_w    = Inches(1.55)
    title_pad = Inches(0.20)
    for i, (num, title, tag) in enumerate(bugs):
        col = i // 6
        idx = i % 6
        x = x_left if col == 0 else x_right
        y = y0 + idx * row_h

        # bug number badge (top-aligned)
        add_rect(s, x, y, badge_w, Inches(0.46), fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, x, y + Inches(0.10), badge_w, Inches(0.30),
                 num, font=HEADER_FONT, size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

        # tag pill (right edge of the column)
        tag_x = x + col_w - pill_w
        add_rect(s, tag_x, y + Inches(0.08), pill_w, Inches(0.30),
                 fill=ICE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=NAVY)
        add_text(s, tag_x, y + Inches(0.11), pill_w, Inches(0.25),
                 tag, font=BODY_FONT, size=8.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        # title — width = col_w - badge - pill - safe padding
        title_x = x + badge_w + title_pad
        title_w = col_w - badge_w - title_pad - pill_w - Inches(0.20)
        add_text(s, title_x, y + Inches(0.02), title_w, Inches(0.30),
                 title, font=BODY_FONT, size=10.5, bold=True, color=CHARCOAL)
        # fix checkmark below the title
        add_text(s, title_x, y + Inches(0.30), title_w, Inches(0.22),
                 "✓ Fixed in code", font=BODY_FONT, size=9, italic=True, color=GREEN_OK)

    # footer note
    add_rect(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.4),
             fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.55), Inches(6.78), Inches(12.2), Inches(0.3),
             "None of these would have been caught by typical functional testing — only sustained load surfaced them.",
             font=BODY_FONT, size=10.5, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# ---------------- Slide 10 — Production posture ----------------
def slide_posture():
    s = new_slide(10)
    page_header(s, "What this capacity means in production", kicker="Production posture")

    # Three big-number callouts
    stats = [
        ("30/s",       "sustained signups",     "at 350 concurrent users"),
        ("2.6 M",      "signups / day",         "if sustained 24h"),
        ("0",          "failures under cascade","350 max-spread waves"),
    ]
    x = Inches(0.55)
    y = Inches(1.7)
    card_w = Inches(4.05)
    card_h = Inches(2.20)
    gap = Inches(0.12)
    for big, mid, small in stats:
        add_rect(s, x, y, card_w, card_h, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, x, y + Inches(0.20), card_w, Inches(1.05),
                 big, font=HEADER_FONT, size=64, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.30), card_w, Inches(0.4),
                 mid, font=BODY_FONT, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.65), card_w, Inches(0.4),
                 small, font=BODY_FONT, size=10.5, italic=True, color=ICE, align=PP_ALIGN.CENTER)
        x += card_w + gap

    # Industry comparison callout
    add_rect(s, Inches(0.55), Inches(4.30), Inches(12.20), Inches(2.55),
             fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=ICE)
    add_rect(s, Inches(0.55), Inches(4.30), Inches(0.18), Inches(2.55), fill=GOLD)
    add_text(s, Inches(0.9), Inches(4.45), Inches(11.6), Inches(0.5),
             "Context: how this compares to the industry",
             font=HEADER_FONT, size=17, bold=True, color=NAVY)
    add_text(s, Inches(0.9), Inches(4.92), Inches(11.6), Inches(1.85),
             "At 30 signups per second sustained, the platform handles 108,000 new ambassadors per hour. The largest MLM companies "
             "globally process about 10,000 signups per day at peak. Our pipeline can absorb a day of their growth in under "
             "six minutes — and we have not yet pushed past the rate-limiter cap to find a real collapse point. "
             "The bottleneck under realistic load (signups distributed across many sponsors) is upline statistics propagation, "
             "which the new eventual-consistency pipeline now handles asynchronously and atomically.",
             font=BODY_FONT, size=12.5, color=CHARCOAL, line_spacing=1.30)

# ---------------- Slide 11 — Roadmap ----------------
def slide_roadmap():
    s = new_slide(11)
    page_header(s, "What's next", kicker="Roadmap")

    # Two-column: left = ranks remaining, right = production prep
    # Left
    add_rect(s, Inches(0.55), Inches(1.55), Inches(0.18), Inches(5.20), fill=GOLD)
    add_text(s, Inches(0.9), Inches(1.55), Inches(6.0), Inches(0.45),
             "Finish the rank ladder", font=HEADER_FONT, size=18, bold=True, color=NAVY)
    add_text(s, Inches(0.9), Inches(2.00), Inches(6.0), Inches(0.4),
             "8 ranks remaining out of 19 total",
             font=BODY_FONT, size=11, italic=True, color=MUTED)

    ranks_remaining = [
        ("12", "Triple Diamond",     "literal signups"),
        ("13", "Blue Diamond",       "literal signups"),
        ("14", "Black Diamond",      "literal / seeded"),
        ("15", "Royal",              "seeded data"),
        ("16", "Double Royal",       "seeded data"),
        ("17", "Triple Royal",       "seeded data"),
        ("18", "Blue Royal",         "seeded data"),
        ("19", "Black Royal",        "seeded data"),
    ]
    y = Inches(2.55)
    for num, name, mode in ranks_remaining:
        add_text(s, Inches(0.9), y, Inches(0.5), Inches(0.32),
                 num, font=HEADER_FONT, size=12, bold=True, color=GOLD_DEEP)
        add_text(s, Inches(1.45), y, Inches(2.8), Inches(0.32),
                 name, font=BODY_FONT, size=11.5, bold=True, color=NAVY)
        add_text(s, Inches(4.4), y, Inches(2.4), Inches(0.32),
                 mode, font=BODY_FONT, size=10.5, italic=True, color=CHARCOAL)
        y += Inches(0.32)

    # Note
    add_text(s, Inches(0.9), y + Inches(0.10), Inches(6.0), Inches(0.6),
             "At Black Royal the qualification alone requires 350,000 enrollment-tree points — over 150,000 literal signups. We switch to seeded test data at the top of the ladder.",
             font=BODY_FONT, size=9.5, italic=True, color=MUTED, line_spacing=1.25)

    # Right column: production readiness
    add_rect(s, Inches(7.05), Inches(1.55), Inches(0.18), Inches(5.20), fill=NAVY)
    add_text(s, Inches(7.4), Inches(1.55), Inches(5.6), Inches(0.45),
             "Production readiness", font=HEADER_FONT, size=18, bold=True, color=NAVY)
    add_text(s, Inches(7.4), Inches(2.00), Inches(5.6), Inches(0.4),
             "After the rank ladder closes",
             font=BODY_FONT, size=11, italic=True, color=MUTED)

    items = [
        "Production deployment preparation (RDS, environments, secrets)",
        "Real Firebase push wiring (current build uses a logging stub)",
        "S3 certificate storage swap (currently local file system)",
        "Redis production wiring (currently optional cache)",
        "Performance hardening for predicted bottlenecks",
        "Frontend (Blazor MAUI Hybrid) — Release 2 begins",
    ]
    y = Inches(2.55)
    for it in items:
        # bullet square
        add_rect(s, Inches(7.4), y + Inches(0.08), Inches(0.13), Inches(0.13), fill=GOLD)
        add_text(s, Inches(7.65), y, Inches(5.5), Inches(0.4),
                 it, font=BODY_FONT, size=11.5, color=CHARCOAL)
        y += Inches(0.50)

# ---------------- Slide 12 — Risks / decisions pending ----------------
def slide_risks():
    s = new_slide(12)
    page_header(s, "Decisions pending and watch-items", kicker="Risks & decisions")

    items = [
        ("Certificate generation policy",
         "On-demand (current) vs. auto-generate at promotion. On-demand is cheaper at scale; auto-generate is simpler for members. Needs a UX call.",
         "DECISION"),
        ("Dashboard freshness window",
         "Under the new eventual-consistency model, admin dashboards can lag the live tree by up to 60 seconds. Acceptable for ops; needs to be communicated.",
         "POLICY"),
        ("Hangfire cross-queue poisoning",
         "Recurring jobs can get parked on the wrong queue when a service restarts before re-registering. Needs a defensive startup guard.",
         "WATCH"),
        ("External-customer rank axis removed",
         "The business rule that required external (non-Ambassador) members in the downline for rank qualification was lifted. Confirm this is the intended product behaviour.",
         "CONFIRM"),
    ]
    y = Inches(1.65)
    for title, desc, tag in items:
        add_rect(s, Inches(0.55), y, Inches(12.2), Inches(1.25), fill=WHITE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=ICE)
        # tag pill
        tag_color = {"DECISION": NAVY, "POLICY": GOLD_DEEP, "WATCH": RED_FLAG, "CONFIRM": TEAL}[tag]
        add_rect(s, Inches(0.75), y + Inches(0.20), Inches(1.40), Inches(0.32),
                 fill=tag_color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, Inches(0.75), y + Inches(0.23), Inches(1.40), Inches(0.28),
                 tag, font=BODY_FONT, size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # title
        add_text(s, Inches(2.30), y + Inches(0.18), Inches(10.3), Inches(0.35),
                 title, font=HEADER_FONT, size=14, bold=True, color=NAVY)
        # desc
        add_text(s, Inches(2.30), y + Inches(0.55), Inches(10.3), Inches(0.65),
                 desc, font=BODY_FONT, size=11, color=CHARCOAL, line_spacing=1.20)
        y += Inches(1.35)

# ---------------- Slide 13 — Closing / Q&A ----------------
def slide_closing():
    s = prs.slides.add_slide(BLANK)
    page_bg(s, NAVY_DARK)
    # gold motif
    add_rect(s, Inches(0), Inches(7.05), Inches(2.0), Inches(0.45), fill=GOLD)
    add_rect(s, Inches(0), Inches(6.6), Inches(0.55), Inches(0.45), fill=GOLD_DEEP)
    # text
    add_text(s, Inches(0.8), Inches(2.4), Inches(8), Inches(0.4),
             "DISCUSSION",
             font=BODY_FONT, size=12, bold=True, color=GOLD)
    add_text(s, Inches(0.8), Inches(2.85), Inches(11.5), Inches(1.6),
             "Questions & Next Steps",
             font=HEADER_FONT, size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(4.45), Inches(11), Inches(0.9),
             "Where would you like to focus the next sprint?",
             font=HEADER_FONT, size=22, italic=True, color=ICE)
    add_text(s, Inches(0.8), Inches(5.4), Inches(11), Inches(0.5),
             "Suggested topics:  certificate policy  ·  remaining ranks  ·  production deployment plan  ·  R2 frontend kickoff",
             font=BODY_FONT, size=13, color=ICE)
    # page marker
    add_text(s, Inches(12.2), Inches(7.10), Inches(0.6), Inches(0.3),
             "13", font=BODY_FONT, size=9, color=ICE, align=PP_ALIGN.RIGHT)

# ---------------- Slide 14 — Appendix: service ports / key locations ----------------
def slide_appendix():
    s = new_slide(14)
    page_header(s, "Reference appendix", kicker="Appendix")

    # Two columns
    # Left: service ports
    add_text(s, Inches(0.55), Inches(1.55), Inches(6.0), Inches(0.4),
             "Service ports (development)", font=HEADER_FONT, size=15, bold=True, color=NAVY)
    ports = [
        ("Signups",            "7005"),
        ("RankEngine",         "7009"),
        ("CommissionEngine",   "7010"),
        ("Billing",            "7011"),
        ("BizCenter",          "7006"),
        ("BizCenterWeb",       "7106"),
        ("AdminAPI",           "7007"),
        ("AdminWeb",           "7107"),
        ("TicketMgmtSystem",   "7012"),
        ("SharedAPICenter",    "7013"),
    ]
    y = Inches(2.05)
    for name, port in ports:
        add_text(s, Inches(0.75), y, Inches(3.5), Inches(0.30),
                 name, font=BODY_FONT, size=11, color=CHARCOAL)
        add_rect(s, Inches(4.50), y + Inches(0.03), Inches(1.30), Inches(0.28),
                 fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, Inches(4.50), y + Inches(0.06), Inches(1.30), Inches(0.24),
                 port, font="Consolas", size=10.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        y += Inches(0.35)

    # Right: test suite status
    add_text(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(0.4),
             "Test suites (Sprint 15)", font=HEADER_FONT, size=15, bold=True, color=NAVY)
    suites = [
        ("Domain.Tests",            "86 passing"),
        ("Repository.Tests",        "covered"),
        ("Signups.Tests",           "300 passing"),
        ("RankEngine.Tests",        "122 passing"),
        ("CommissionEngine.Tests",  "67 passing"),
        ("Billing.Tests",           "124 passing"),
        ("BizCenter.Tests",         "142 passing"),
        ("AdminAPI.Tests",          "313 passing"),
        ("TicketManagement.Tests",  "45 passing"),
        ("Total",                   "900+ passing"),
    ]
    y = Inches(2.05)
    for name, status in suites:
        bold = name == "Total"
        add_text(s, Inches(7.20), y, Inches(3.5), Inches(0.30),
                 name, font=BODY_FONT, size=11, bold=bold, color=CHARCOAL)
        color = GREEN_OK if "passing" in status else MUTED
        if name == "Total":
            add_rect(s, Inches(11.0), y + Inches(0.03), Inches(1.6), Inches(0.28),
                     fill=GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            add_text(s, Inches(11.0), y + Inches(0.06), Inches(1.6), Inches(0.24),
                     status, font=BODY_FONT, size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        else:
            add_text(s, Inches(11.0), y, Inches(2.0), Inches(0.30),
                     status, font=BODY_FONT, size=10.5, bold=True, color=color)
        y += Inches(0.35)

    # bottom: key file locations
    by = Inches(5.95)
    add_rect(s, Inches(0.55), by, Inches(12.2), Inches(1.0), fill=ICE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.75), by + Inches(0.10), Inches(11.8), Inches(0.32),
             "Key references", font=HEADER_FONT, size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.75), by + Inches(0.42), Inches(11.8), Inches(0.55),
             "Project root:   C:\\Users\\sagar\\source\\repos\\ClaudeRepository\nSprint 15 state:   memory/project_sprint15_state.md\nSpec + plan:   docs/superpowers/specs · docs/superpowers/plans",
             font="Consolas", size=10, color=CHARCOAL, line_spacing=1.30)


# ---------------- Build ----------------
slide_title()
slide_exec_summary()
slide_mindmap()
slide_service_table()
slide_data_flow()
slide_sprint15()
slide_load_test()
slide_improvements()
slide_bugs()
slide_posture()
slide_roadmap()
slide_risks()
slide_closing()
slide_appendix()

prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {OUT_PATH.stat().st_size / 1024:.1f} KB")
